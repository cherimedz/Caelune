import html as _html
import io
import os
import re as _re
from datetime import datetime

import markdown as _md
import streamlit as st
from huggingface_hub import InferenceClient
from prompts import (
    build_messages, build_key_terms_messages, build_quiz_messages, build_chat_messages,
    max_tokens_for,
    FORMAT_LABELS, AVAILABLE_MODELS, DEFAULT_MODEL, VISION_MODEL,
    SUBJECTS, DIFFICULTIES, LENGTHS, LANGUAGES,
)

st.set_page_config(
    page_title="Caelune",
    page_icon="assets/icon.png" if os.path.exists("assets/icon.png") else None,
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── External resources ────────────────────────────────────────────────────────
st.markdown(
    '<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.5.1/css/all.min.css"/>',
    unsafe_allow_html=True,
)
st.markdown(
    '<link href="https://fonts.googleapis.com/css2?family=Playfair+Display:ital,wght@0,400;0,600;0,700;1,400&family=DM+Sans:ital,wght@0,300;0,400;0,500;0,600;1,400&display=swap" rel="stylesheet"/>',
    unsafe_allow_html=True,
)

# ── CSS ───────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
*, *::before, *::after { box-sizing: border-box; }

/* ── Base ── */
html, body,
.stApp,
[data-testid="stAppViewContainer"],
[data-testid="stMain"] {
  font-family: 'DM Sans', system-ui, sans-serif !important;
  background: transparent !important;
  color: #F9CCCB !important;
}

/* ── Fixed twilight background ── */
.stApp::before {
  content: '';
  position: fixed;
  inset: 0;
  z-index: -10;
  background: linear-gradient(
    170deg,
    #082D36 0%,
    #197074 18%,
    #2a0c12 36%,
    #450A18 52%,
    #7B182F 65%,
    #C43556 80%,
    #F9CCCB 100%
  );
}

/* ── Chrome cleanup ── */
[data-testid="stHeader"]  { background: transparent !important; }
[data-testid="stToolbar"] { display: none !important; }
#MainMenu, footer          { visibility: hidden !important; }

/* ── Sidebar — dark glass ── */
section[data-testid="stSidebar"] {
  background: rgba(8,45,54,0.92) !important;
  backdrop-filter: blur(20px) saturate(1.4) !important;
  border-right: 2px solid rgba(196,53,86,0.35) !important;
}
section[data-testid="stSidebar"] * {
  font-family: 'DM Sans', sans-serif !important;
  color: rgba(249,204,203,0.9) !important;
  -webkit-text-fill-color: rgba(249,204,203,0.9) !important;
}

/* ── Sidebar collapse button (the arrow inside open sidebar) ── */
[data-testid="stSidebarCollapseButton"] button {
  background: rgba(196,53,86,0.18) !important;
  border: 1px solid rgba(249,204,203,0.35) !important;
  border-radius: 50% !important;
  color: #F9CCCB !important;
  width: 34px !important;
  height: 34px !important;
  display: flex !important;
  align-items: center !important;
  justify-content: center !important;
  box-shadow: 0 0 14px rgba(196,53,86,0.25) !important;
  transition: all 0.2s !important;
}
[data-testid="stSidebarCollapseButton"] button:hover {
  background: rgba(196,53,86,0.35) !important;
  box-shadow: 0 0 22px rgba(196,53,86,0.5) !important;
  transform: scale(1.08) !important;
}
[data-testid="stSidebarCollapseButton"] button svg {
  fill: #F9CCCB !important;
  stroke: #F9CCCB !important;
}

/* ── Collapsed sidebar control (re-open pill after sidebar is closed) ── */
[data-testid="stSidebarCollapsedControl"] {
  background: rgba(8,45,54,0.92) !important;
  backdrop-filter: blur(16px) !important;
  border-radius: 0 20px 20px 0 !important;
  border-right: 2px solid rgba(196,53,86,0.5) !important;
  border-top: 1px solid rgba(134,221,228,0.12) !important;
  border-bottom: 1px solid rgba(134,221,228,0.12) !important;
  box-shadow: 4px 0 28px rgba(196,53,86,0.3) !important;
  padding: 10px 6px !important;
  animation: pillPulse 3s ease-in-out infinite !important;
}
@keyframes pillPulse {
  0%,100% { box-shadow: 4px 0 24px rgba(196,53,86,0.25); }
  50%      { box-shadow: 4px 0 36px rgba(196,53,86,0.55); }
}
[data-testid="stSidebarCollapsedControl"] button {
  background: rgba(196,53,86,0.22) !important;
  border: 1px solid rgba(249,204,203,0.35) !important;
  border-radius: 50% !important;
  color: #F9CCCB !important;
  width: 30px !important;
  height: 30px !important;
  display: flex !important;
  align-items: center !important;
  justify-content: center !important;
  transition: all 0.2s !important;
}
[data-testid="stSidebarCollapsedControl"] button:hover {
  background: rgba(196,53,86,0.45) !important;
  transform: scale(1.1) !important;
}
[data-testid="stSidebarCollapsedControl"] button svg {
  fill: #F9CCCB !important;
  stroke: #F9CCCB !important;
}

/* ── Main block container — glass card ── */
.block-container {
  background: rgba(8,45,54,0.35) !important;
  backdrop-filter: blur(28px) saturate(1.4) !important;
  border: 1px solid rgba(134,221,228,0.14) !important;
  border-radius: 24px !important;
  box-shadow:
    0 32px 80px rgba(0,0,0,0.5),
    inset 0 1px 0 rgba(134,221,228,0.06) !important;
  padding: 2.5rem 2.8rem !important;
  margin-top: 1.8rem !important;
  margin-bottom: 2rem !important;
}

/* ── Headings inside main ── */
h1 {
  font-family: 'Playfair Display', serif !important;
  font-weight: 700 !important;
  font-size: 2.6rem !important;
  letter-spacing: -0.5px !important;
  color: #F9CCCB !important;
  -webkit-text-fill-color: #F9CCCB !important;
  text-align: center !important;
  margin-bottom: 2px !important;
}
h2, h3 {
  font-family: 'DM Sans', sans-serif !important;
  font-weight: 600 !important;
  font-size: 0.72rem !important;
  letter-spacing: 1.5px !important;
  text-transform: uppercase !important;
  color: rgba(134,221,228,0.65) !important;
  -webkit-text-fill-color: rgba(134,221,228,0.65) !important;
}

/* ── Text area ── */
.stTextArea textarea {
  font-family: 'DM Sans', sans-serif !important;
  font-size: 0.92rem !important;
  background: rgba(249,230,228,0.96) !important;
  border: 1px solid rgba(196,53,86,0.22) !important;
  border-radius: 14px !important;
  color: #450A18 !important;
  line-height: 1.75 !important;
  padding: 14px 16px !important;
  transition: border-color 0.2s, box-shadow 0.2s !important;
  caret-color: #C43556 !important;
  resize: vertical !important;
}
.stTextArea textarea:focus {
  border-color: rgba(196,53,86,0.55) !important;
  box-shadow: 0 0 0 3px rgba(196,53,86,0.12) !important;
  outline: none !important;
}
.stTextArea textarea::placeholder {
  color: rgba(69,10,24,0.4) !important;
  font-style: italic !important;
}

/* ── Select box ── */
[data-baseweb="select"] > div {
  background: rgba(8,45,54,0.4) !important;
  border: 1px solid rgba(134,221,228,0.18) !important;
  border-radius: 10px !important;
  color: #F9CCCB !important;
}
[data-baseweb="select"] * { color: #F9CCCB !important; font-family: 'DM Sans', sans-serif !important; }
[data-baseweb="popover"] { background: #082D36 !important; border: 1px solid rgba(134,221,228,0.15) !important; }
[data-baseweb="menu"] { background: #082D36 !important; }
[data-baseweb="option"] { color: #F9CCCB !important; }
[data-baseweb="option"]:hover { background: rgba(134,221,228,0.08) !important; }

/* ── Radio ── */
div[data-testid="stRadio"] > div { flex-direction: column !important; gap: 6px !important; }
div[data-testid="stRadio"] label {
  background: rgba(8,45,54,0.3) !important;
  border: 1px solid rgba(134,221,228,0.12) !important;
  border-radius: 10px !important;
  padding: 9px 14px !important;
  color: rgba(249,204,203,0.75) !important;
  -webkit-text-fill-color: rgba(249,204,203,0.75) !important;
  font-weight: 500 !important;
  font-size: 0.87rem !important;
  transition: all 0.2s !important;
  cursor: pointer !important;
  width: 100% !important;
}
div[data-testid="stRadio"] label:has(input:checked) {
  background: rgba(196,53,86,0.18) !important;
  border-color: rgba(196,53,86,0.5) !important;
  color: #F9CCCB !important;
  -webkit-text-fill-color: #F9CCCB !important;
}
div[data-testid="stRadio"] input[type="radio"] { display: none !important; }

/* ── Primary button ── */
.stButton > button[kind="primary"] {
  background: linear-gradient(135deg, #C43556 0%, #7B182F 100%) !important;
  border: none !important;
  border-radius: 12px !important;
  font-family: 'DM Sans', sans-serif !important;
  font-weight: 600 !important;
  font-size: 0.93rem !important;
  padding: 0.65rem 1.8rem !important;
  color: #fff !important;
  -webkit-text-fill-color: #fff !important;
  letter-spacing: 0.2px !important;
  box-shadow: 0 4px 20px rgba(196,53,86,0.4) !important;
  transition: transform 0.15s, box-shadow 0.15s, opacity 0.15s !important;
  width: 100% !important;
}
.stButton > button[kind="primary"]:hover:not(:disabled) {
  transform: translateY(-2px) !important;
  box-shadow: 0 8px 30px rgba(196,53,86,0.55) !important;
}
.stButton > button[kind="primary"]:active { transform: translateY(0) scale(0.98) !important; }
.stButton > button[kind="primary"]:disabled { opacity: 0.38 !important; }

/* ── Secondary button (format toggles) ── */
.stButton > button[kind="secondary"] {
  border-radius: 10px !important;
  font-family: 'DM Sans', sans-serif !important;
  font-weight: 500 !important;
  font-size: 0.88rem !important;
  border: 1px solid rgba(134,221,228,0.15) !important;
  color: rgba(249,204,203,0.75) !important;
  -webkit-text-fill-color: rgba(249,204,203,0.75) !important;
  background: rgba(8,45,54,0.25) !important;
  transition: all 0.2s !important;
  width: 100% !important;
}
.stButton > button[kind="secondary"]:hover {
  background: rgba(134,221,228,0.08) !important;
  border-color: rgba(134,221,228,0.3) !important;
}

/* ── Download button ── */
.stDownloadButton > button {
  border-radius: 10px !important;
  font-family: 'DM Sans', sans-serif !important;
  font-weight: 500 !important;
  font-size: 0.88rem !important;
  background: rgba(8,45,54,0.4) !important;
  border: 1px solid rgba(134,221,228,0.2) !important;
  color: #F9CCCB !important;
  -webkit-text-fill-color: #F9CCCB !important;
  transition: all 0.2s !important;
  width: 100% !important;
}
.stDownloadButton > button:hover {
  background: rgba(134,221,228,0.1) !important;
  border-color: rgba(134,221,228,0.35) !important;
}

/* ── Output card ── */
.output-card {
  background: rgba(8,45,54,0.3);
  border: 1px solid rgba(134,221,228,0.12);
  border-radius: 14px;
  padding: 1.3rem 1.5rem;
  min-height: 440px;
  font-family: 'DM Sans', sans-serif;
  font-size: 0.92rem;
  line-height: 1.85;
  color: #F9CCCB;
  white-space: pre-wrap;
  word-break: break-word;
  overflow-y: auto;
}
.output-card.empty {
  display: flex; flex-direction: column;
  align-items: center; justify-content: center;
  gap: 12px;
  color: rgba(249,204,203,0.28);
  font-style: italic; font-size: 0.9rem;
}
.output-card.empty i { font-size: 1.8rem; color: rgba(196,53,86,0.4); }

/* ── Markdown-rendered output card ── */
.output-card.md-output { white-space: normal !important; }
.output-card.md-output h2, .output-card.md-output h3 {
  font-family: 'DM Sans', sans-serif !important;
  font-size: 0.78rem !important;
  font-weight: 700 !important;
  letter-spacing: 1.4px !important;
  text-transform: uppercase !important;
  color: #86DDE4 !important;
  -webkit-text-fill-color: #86DDE4 !important;
  margin: 1.1rem 0 0.4rem !important;
  border-bottom: 1px solid rgba(134,221,228,0.15) !important;
  padding-bottom: 4px !important;
}
.output-card.md-output strong { color: #F9CCCB; font-weight: 600; }
.output-card.md-output em { color: rgba(249,204,203,0.75); }
.output-card.md-output ul, .output-card.md-output ol {
  padding-left: 1.4rem; margin: 0.4rem 0;
}
.output-card.md-output li { margin-bottom: 0.4rem; }
.output-card.md-output p { margin: 0.45rem 0; }
.output-card.md-output table {
  border-collapse: collapse; width: 100%; margin: 0.8rem 0;
}
.output-card.md-output th, .output-card.md-output td {
  border: 1px solid rgba(134,221,228,0.2);
  padding: 6px 10px; font-size: 0.88rem;
}
.output-card.md-output th {
  background: rgba(134,221,228,0.08); font-weight: 600;
}

/* ── Section header labels ── */
.sect-label {
  display: flex; align-items: center; gap: 8px;
  font-size: 0.72rem; font-weight: 600;
  text-transform: uppercase; letter-spacing: 1.5px;
  color: rgba(134,221,228,0.65);
  margin-bottom: 8px;
}
.sect-label i { font-size: 0.8rem; color: rgba(134,221,228,0.9); }

/* ── Sidebar section labels ── */
.sidebar-label {
  font-size: 0.7rem; font-weight: 600;
  text-transform: uppercase; letter-spacing: 1.4px;
  color: rgba(134,221,228,0.55);
  margin-bottom: 8px; display: flex; align-items: center; gap: 7px;
}
.sidebar-label i { color: rgba(134,221,228,0.8); }

/* ── Info rows (sidebar) ── */
.info-row {
  display: flex; align-items: center; gap: 9px;
  font-size: 0.8rem; color: rgba(249,204,203,0.6);
  padding: 5px 0;
}
.info-row i { width: 14px; color: rgba(134,221,228,0.75); font-size: 0.78rem; }

/* ── Caption ── */
.stCaption p {
  color: rgba(249,204,203,0.38) !important;
  font-family: 'DM Sans', sans-serif !important;
  font-size: 0.75rem !important;
}

/* ── Alerts ── */
[data-testid="stAlert"] {
  border-radius: 12px !important;
  font-family: 'DM Sans', sans-serif !important;
  font-size: 0.88rem !important;
}

/* ── Divider ── */
hr { border-color: rgba(134,221,228,0.1) !important; margin: 1rem 0 !important; }

/* ── Blinking cursor for streaming ── */
@keyframes blink { 0%,100%{opacity:1} 50%{opacity:0} }
.cursor { animation: blink 1s step-end infinite; color: #C43556; }

/* ── Hide Streamlit's keyboard-shortcut hint text ── */
[data-testid*="keyboard"],
[data-testid="stSidebarHeader"],
[class*="shortcut"],
kbd { display: none !important; }

/* ── Quiz ── */
.quiz-q {
  background: rgba(8,45,54,0.25);
  border: 1px solid rgba(134,221,228,0.12);
  border-radius: 12px;
  padding: 0.85rem 1.05rem;
  margin-bottom: 0.75rem;
}
.quiz-q.quiz-correct {
  border-color: rgba(25,112,116,0.55) !important;
  background: rgba(25,112,116,0.1) !important;
}
.quiz-q.quiz-wrong {
  border-color: rgba(196,53,86,0.5) !important;
  background: rgba(196,53,86,0.07) !important;
}
.quiz-q-text { color: #F9CCCB; font-weight: 600; font-size: 0.9rem; }
.quiz-badge {
  display: inline-block; font-size: 0.67rem; font-weight: 700;
  letter-spacing: 0.8px; text-transform: uppercase;
  padding: 2px 9px; border-radius: 20px; margin-bottom: 0.4rem;
}
.quiz-badge.correct { background: rgba(25,112,116,0.25); color: #86DDE4; }
.quiz-badge.wrong   { background: rgba(196,53,86,0.2);   color: #F9CCCB; }
.quiz-score {
  text-align: center; padding: 1.1rem 1rem;
  background: rgba(8,45,54,0.35); border-radius: 14px;
  border: 1px solid rgba(134,221,228,0.15); margin-bottom: 1rem;
}
.quiz-score-num {
  font-family: 'Playfair Display', serif; font-size: 2rem;
  font-weight: 700; color: #86DDE4;
}
.quiz-score-label { font-size: 0.78rem; color: rgba(249,204,203,0.5); letter-spacing: 1.1px; }

/* ── Chat ── */
.chat-container {
  display: flex; flex-direction: column; gap: 10px;
  max-height: 360px; overflow-y: auto; padding: 0.4rem 0 0.6rem;
}
.chat-bubble {
  padding: 0.7rem 1rem; border-radius: 14px;
  font-size: 0.9rem; line-height: 1.65; max-width: 88%;
  word-break: break-word;
}
.chat-bubble.user {
  background: rgba(196,53,86,0.18); border: 1px solid rgba(196,53,86,0.28);
  color: #F9CCCB; margin-left: auto; border-bottom-right-radius: 4px;
}
.chat-bubble.ai {
  background: rgba(8,45,54,0.4); border: 1px solid rgba(134,221,228,0.12);
  color: #F9CCCB; border-bottom-left-radius: 4px;
}
.chat-bubble.ai p { margin: 0.3rem 0; }
.chat-bubble.ai strong { color: #F9CCCB; font-weight: 600; }
.chat-bubble.ai ul, .chat-bubble.ai ol { padding-left: 1.2rem; margin: 0.3rem 0; }
.chat-bubble.ai h2, .chat-bubble.ai h3 {
  font-size: 0.74rem !important; text-transform: uppercase;
  color: #86DDE4 !important; -webkit-text-fill-color: #86DDE4 !important;
  letter-spacing: 1.1px; margin: 0.5rem 0 0.2rem;
}
.chat-who {
  font-size: 0.62rem; font-weight: 700; letter-spacing: 1.2px;
  text-transform: uppercase; opacity: 0.5; margin-bottom: 4px;
}
</style>
""", unsafe_allow_html=True)

# ── Background scene ─────────────────────────────────────────────────────────
st.markdown("""
<div style="position:fixed;inset:0;z-index:-5;pointer-events:none;overflow:hidden;">

  <!-- Sky canvas: stars, aurora, shooting stars, particles, constellations -->
  <canvas id="skyCanvas" style="position:absolute;inset:0;width:100%;height:100%;"></canvas>

  <!-- Moon -->
  <div style="
    position:absolute; top:6%; right:9%;
    width:64px; height:64px; border-radius:50%;
    background:radial-gradient(circle at 38% 32%, #ffffff, #F9CCCB 60%, #86DDE4);
    box-shadow:
      0 0 0 2px rgba(134,221,228,0.15),
      0 0 30px 10px rgba(249,204,203,0.4),
      0 0 80px 30px rgba(134,221,228,0.18),
      0 0 140px 60px rgba(25,112,116,0.08);
    overflow:hidden;">
    <div style="position:absolute;top:-6px;left:13px;width:58px;height:58px;border-radius:50%;background:#082D36;opacity:0.75;"></div>
  </div>

  <!-- Ringed planet (upper left) -->
  <svg style="position:absolute;top:9%;left:7%;opacity:0.75;" width="54" height="36" viewBox="0 0 54 36">
    <defs>
      <radialGradient id="pg" cx="50%" cy="40%">
        <stop offset="0%"   stop-color="#86DDE4"/>
        <stop offset="100%" stop-color="#197074"/>
      </radialGradient>
    </defs>
    <!-- ring back half -->
    <ellipse cx="27" cy="22" rx="26" ry="6" fill="none" stroke="rgba(134,221,228,0.35)" stroke-width="3"/>
    <!-- planet body -->
    <circle cx="27" cy="18" r="11" fill="url(#pg)" opacity="0.9"/>
    <!-- planet glow -->
    <circle cx="27" cy="18" r="11" fill="none" stroke="rgba(134,221,228,0.25)" stroke-width="4"/>
    <!-- ring front half -->
    <path d="M1,22 Q27,30 53,22" fill="none" stroke="rgba(134,221,228,0.45)" stroke-width="2.5"/>
  </svg>

  <!-- Mountain layers -->
  <svg viewBox="0 0 1440 360" preserveAspectRatio="xMidYMax slice"
       style="position:absolute;bottom:0;left:0;width:100%;height:55%;">
    <path d="M0,280 L90,200 L180,240 L290,150 L400,210 L520,130 L640,185 L760,110 L880,170 L1000,105 L1110,160 L1230,130 L1340,175 L1440,150 L1440,360 L0,360Z"
          fill="rgba(8,45,54,0.55)"/>
    <path d="M0,310 L70,265 L160,290 L260,240 L370,270 L480,225 L590,255 L700,210 L810,245 L920,215 L1030,250 L1150,225 L1270,255 L1380,240 L1440,260 L1440,360 L0,360Z"
          fill="rgba(4,22,27,0.80)"/>
    <path d="M0,340 L80,315 L160,330 L260,305 L360,325 L460,300 L560,318 L660,295 L760,312 L860,290 L960,308 L1060,285 L1160,305 L1260,288 L1360,308 L1440,295 L1440,360 L0,360Z"
          fill="rgba(2,10,14,0.92)"/>
    <!-- trees left -->
    <rect x="18" y="330" width="7" height="28" fill="#020a0e"/>
    <polygon points="21.5,302 11,333 32,333" fill="#020a0e"/>
    <polygon points="21.5,312 9,338 34,338" fill="#020a0e"/>
    <rect x="45" y="334" width="5" height="24" fill="#020a0e"/>
    <polygon points="47.5,308 39,336 56,336" fill="#020a0e"/>
    <rect x="68" y="337" width="5" height="21" fill="#020a0e"/>
    <polygon points="70.5,314 62,339 79,339" fill="#020a0e"/>
    <!-- trees right -->
    <rect x="1358" y="328" width="7" height="30" fill="#020a0e"/>
    <polygon points="1361.5,300 1351,331 1372,331" fill="#020a0e"/>
    <polygon points="1361.5,310 1349,336 1374,336" fill="#020a0e"/>
    <rect x="1385" y="332" width="5" height="26" fill="#020a0e"/>
    <polygon points="1387.5,306 1379,335 1396,335" fill="#020a0e"/>
    <rect x="1408" y="335" width="5" height="23" fill="#020a0e"/>
    <polygon points="1410.5,312 1402,338 1419,338" fill="#020a0e"/>
  </svg>
</div>

<script>
(function(){
  function initSky(){
    const canvas = document.getElementById('skyCanvas');
    if(!canvas){ setTimeout(initSky,200); return; }
    const ctx = canvas.getContext('2d');
    const W = ()=>canvas.width, H = ()=>canvas.height;

    function resize(){
      canvas.width  = window.innerWidth;
      canvas.height = window.innerHeight;
    }
    resize();
    window.addEventListener('resize', ()=>{ resize(); initStars(); initParticles(); });

    /* ── Stars ── */
    let stars = [];
    function initStars(){
      stars = Array.from({length:200}, ()=>({
        x: Math.random()*W(),
        y: Math.random()*H()*0.63,
        r: Math.random()*1.7+0.25,
        phase: Math.random()*Math.PI*2,
        speed: 0.002+Math.random()*0.007,
        col: [[255,255,255],[249,204,203],[134,221,228],[255,235,200]][Math.floor(Math.random()*4)]
      }));
    }
    initStars();

    /* ── Constellations ── */
    const constellations = [
      { pts:[{x:.08,y:.13},{x:.11,y:.17},{x:.13,y:.21},{x:.10,y:.26},{x:.14,y:.26},{x:.07,y:.30},{x:.16,y:.30}],
        lines:[[0,1],[1,2],[2,3],[2,4],[3,5],[4,6]] },
      { pts:[{x:.76,y:.07},{x:.80,y:.10},{x:.84,y:.13},{x:.86,y:.18},{x:.83,y:.23}],
        lines:[[0,1],[1,2],[2,3],[3,4]] },
      { pts:[{x:.45,y:.05},{x:.48,y:.09},{x:.52,y:.07},{x:.55,y:.11},{x:.50,y:.14}],
        lines:[[0,1],[2,3],[1,4],[3,4]] },
    ];

    /* ── Shooting stars ── */
    let shooters = [];
    function spawnShooter(delayMs){
      setTimeout(()=>{
        if(shooters.filter(s=>s.alive).length >= 9) return;
        const ang = (18+Math.random()*28)*Math.PI/180;
        const spd = 9+Math.random()*14;
        shooters.push({
          x: Math.random()*W()*0.72,
          y: Math.random()*H()*0.38,
          dx: Math.cos(ang)*spd, dy: Math.sin(ang)*spd,
          len: 70+Math.random()*120,
          w: 1+Math.random()*1.8,
          life:0, maxLife:35+Math.random()*35,
          alive:true,
          col: [[255,255,255],[249,204,203],[134,221,228]][Math.floor(Math.random()*3)]
        });
      }, delayMs||0);
    }

    /* ── Floating particles ── */
    let particles=[];
    const PCOLS=[[249,204,203],[134,221,228],[196,53,86],[25,112,116],[255,255,255]];
    function initParticles(){
      particles = Array.from({length:32},()=>newParticle(true));
    }
    function newParticle(randomY){
      return {
        x: Math.random()*W(),
        y: randomY ? H()*(0.35+Math.random()*0.6) : H()+5,
        vy: -(0.12+Math.random()*0.38),
        vx: (Math.random()-.5)*0.12,
        r: 0.7+Math.random()*1.8,
        life:randomY?Math.floor(Math.random()*200):0,
        maxLife:140+Math.random()*260,
        col:PCOLS[Math.floor(Math.random()*PCOLS.length)]
      };
    }
    initParticles();

    /* ── Aurora ── */
    const auroraBands = [
      {y:.09, amp:20, freq:.0022, ph:0,   col:[134,221,228], op:.08 },
      {y:.16, amp:26, freq:.0028, ph:2.1, col:[196,53,86],   op:.09 },
      {y:.06, amp:15, freq:.0018, ph:3.8, col:[249,204,203], op:.065},
      {y:.21, amp:18, freq:.0032, ph:1.2, col:[25,112,116],  op:.06 },
    ];

    let t=0, lastSpawn=0;

    function drawAurora(){
      auroraBands.forEach(b=>{
        const yBase = b.y*H();
        const [r,g,bl]=b.col;
        const grd = ctx.createLinearGradient(0,0,W(),0);
        grd.addColorStop(0,   `rgba(${r},${g},${bl},0)`);
        grd.addColorStop(.12, `rgba(${r},${g},${bl},${b.op})`);
        grd.addColorStop(.5,  `rgba(${r},${g},${bl},${b.op*1.6})`);
        grd.addColorStop(.88, `rgba(${r},${g},${bl},${b.op})`);
        grd.addColorStop(1,   `rgba(${r},${g},${bl},0)`);
        ctx.save();
        ctx.beginPath();
        ctx.moveTo(0,yBase);
        for(let x=0;x<=W();x+=3){
          const y=yBase
            +Math.sin(x*b.freq+t*.006+b.ph)*b.amp
            +Math.sin(x*b.freq*1.8+t*.004+b.ph+1.3)*(b.amp*.35);
          ctx.lineTo(x,y);
        }
        ctx.lineTo(W(),H()*.5); ctx.lineTo(0,H()*.5);
        ctx.closePath();
        ctx.fillStyle=grd; ctx.filter='blur(6px)';
        ctx.fill(); ctx.restore();
      });
    }

    function drawStars(){
      stars.forEach(s=>{
        const a=.28+.72*Math.abs(Math.sin(t*s.speed+s.phase));
        const [r,g,b]=s.col;
        if(s.r>1.1){
          ctx.beginPath(); ctx.arc(s.x,s.y,s.r*3,0,Math.PI*2);
          ctx.fillStyle=`rgba(${r},${g},${b},${a*.1})`; ctx.fill();
        }
        ctx.beginPath(); ctx.arc(s.x,s.y,s.r,0,Math.PI*2);
        ctx.fillStyle=`rgba(${r},${g},${b},${a})`; ctx.fill();
      });
    }

    function drawConstellations(){
      constellations.forEach(c=>{
        const pulse=.5+.5*Math.sin(t*.003);
        c.lines.forEach(([a,b])=>{
          ctx.beginPath();
          ctx.moveTo(c.pts[a].x*W(), c.pts[a].y*H());
          ctx.lineTo(c.pts[b].x*W(), c.pts[b].y*H());
          ctx.strokeStyle=`rgba(134,221,228,${.08+pulse*.06})`;
          ctx.lineWidth=.6; ctx.stroke();
        });
        c.pts.forEach(p=>{
          const a=.55+.45*Math.abs(Math.sin(t*.004+p.x*8));
          ctx.beginPath(); ctx.arc(p.x*W(),p.y*H(),1.8,0,Math.PI*2);
          ctx.fillStyle=`rgba(249,204,203,${a})`; ctx.fill();
          ctx.beginPath(); ctx.arc(p.x*W(),p.y*H(),5,0,Math.PI*2);
          ctx.fillStyle=`rgba(134,221,228,${a*.14})`; ctx.fill();
        });
      });
    }

    function drawShooters(){
      shooters = shooters.filter(s=>s.alive);
      shooters.forEach(s=>{
        s.life++;
        const prog=s.life/s.maxLife;
        const op = prog<.15 ? prog/.15 : prog>.7 ? 1-(prog-.7)/.3 : 1;
        if(s.life>s.maxLife){ s.alive=false; return; }
        const angle=Math.atan2(s.dy,s.dx);
        const tx=s.x-Math.cos(angle)*s.len, ty=s.y-Math.sin(angle)*s.len;
        const grd=ctx.createLinearGradient(tx,ty,s.x,s.y);
        const [r,g,b]=s.col;
        grd.addColorStop(0,`rgba(${r},${g},${b},0)`);
        grd.addColorStop(.6,`rgba(${r},${g},${b},${op*.45})`);
        grd.addColorStop(1,`rgba(${r},${g},${b},${op})`);
        ctx.beginPath(); ctx.moveTo(tx,ty); ctx.lineTo(s.x,s.y);
        ctx.strokeStyle=grd; ctx.lineWidth=s.w; ctx.lineCap='round'; ctx.stroke();
        /* head glow */
        ctx.beginPath(); ctx.arc(s.x,s.y,s.w*2,0,Math.PI*2);
        ctx.fillStyle=`rgba(255,255,255,${op*.7})`; ctx.fill();
        s.x+=s.dx; s.y+=s.dy;
      });
    }

    function drawParticles(){
      particles.forEach((p,i)=>{
        p.life++; p.x+=p.vx; p.y+=p.vy;
        if(p.life>p.maxLife){ particles[i]=newParticle(false); return; }
        const prog=p.life/p.maxLife;
        const a=prog<.18 ? prog/.18 : prog>.78 ? 1-(prog-.78)/.22 : 1;
        const [r,g,b]=p.col;
        ctx.beginPath(); ctx.arc(p.x,p.y,p.r*3.5,0,Math.PI*2);
        ctx.fillStyle=`rgba(${r},${g},${b},${a*.07})`; ctx.fill();
        ctx.beginPath(); ctx.arc(p.x,p.y,p.r,0,Math.PI*2);
        ctx.fillStyle=`rgba(${r},${g},${b},${a*.5})`; ctx.fill();
      });
    }

    function loop(){
      ctx.clearRect(0,0,W(),H());
      drawAurora();
      drawStars();
      drawConstellations();
      drawShooters();
      drawParticles();

      /* spawn shooting stars */
      if(t-lastSpawn > 22+Math.random()*38){
        spawnShooter(0);
        if(Math.random()<.55) spawnShooter(180+Math.random()*300);
        if(Math.random()<.30) spawnShooter(420+Math.random()*400);
        if(Math.random()<.12) spawnShooter(750+Math.random()*500);
        lastSpawn=t;
      }
      t++;
      requestAnimationFrame(loop);
    }
    loop();
  }
  initSky();
})();
</script>

""", unsafe_allow_html=True)


# ── Helpers ───────────────────────────────────────────────────────────────────
def get_hf_token() -> str | None:
    try:
        return st.secrets["HF_TOKEN"]
    except (KeyError, FileNotFoundError):
        return os.environ.get("HF_TOKEN")


def extract_text_from_file(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    data = io.BytesIO(uploaded_file.read())
    try:
        if name.endswith(".pdf"):
            import pdfplumber
            with pdfplumber.open(data) as pdf:
                pages = [p.extract_text() or "" for p in pdf.pages]
            return "\n\n".join(p for p in pages if p.strip())
        if name.endswith(".docx"):
            from docx import Document
            doc = Document(data)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(data)
            texts = [
                shape.text.strip()
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            return "\n".join(texts)
        if name.endswith(".txt"):
            return data.read().decode("utf-8", errors="ignore")
    except Exception as e:
        st.error(f"Could not read file: {e}")
    return ""


_IMAGE_MIME = {".jpg": "image/jpeg", ".jpeg": "image/jpeg",
               ".png": "image/png", ".webp": "image/webp"}
_IMAGE_EXTS = tuple(_IMAGE_MIME)


def _resize_for_api(raw: bytes, max_dim: int = 1568) -> tuple[bytes, str]:
    from PIL import Image as _Img
    img = _Img.open(io.BytesIO(raw)).convert("RGB")
    if max(img.size) > max_dim:
        img.thumbnail((max_dim, max_dim), _Img.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=88)
    return buf.getvalue(), "image/jpeg"


def extract_text_from_image(raw: bytes, ext: str, token: str) -> str:
    import base64
    resized, mime = _resize_for_api(raw)
    b64 = base64.b64encode(resized).decode()
    client = InferenceClient(api_key=token)
    try:
        resp = client.chat.completions.create(
            model=VISION_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                    {"type": "text", "text": (
                        "You are an OCR engine. Extract ALL text visible in this image exactly "
                        "as written. Preserve paragraph breaks and structure. "
                        "Return only the extracted text — no commentary, no preamble."
                    )},
                ],
            }],
            max_tokens=1800,
            stream=False,
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"Could not extract text from image: {e}")
        return ""


def stream_call(messages: list[dict], token: str, model: str, max_tokens: int):
    client = InferenceClient(api_key=token)
    stream = client.chat.completions.create(
        model=model, messages=messages, max_tokens=max_tokens, stream=True,
    )
    for chunk in stream:
        if not chunk.choices:
            continue
        delta = chunk.choices[0].delta.content
        if delta:
            yield delta


def parse_quiz(text: str) -> list[dict]:
    questions = []
    blocks = _re.split(r'\n(?=Q\d+[\.:]\s)', text.strip())
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        lines = [l.strip() for l in block.splitlines() if l.strip()]
        if not lines:
            continue
        q_match = _re.match(r'^Q\d+[\.:]\s*(.*)', lines[0])
        if not q_match:
            continue
        question = q_match.group(1).strip()
        options: dict[str, str] = {}
        answer = None
        for line in lines[1:]:
            opt_m = _re.match(r'^([A-D])[).]?\s+(.*)', line)
            if opt_m:
                options[opt_m.group(1)] = opt_m.group(2).strip()
                continue
            ans_m = _re.match(r'^Answer:\s*([A-D])', line, _re.IGNORECASE)
            if ans_m:
                answer = ans_m.group(1).upper()
        if question and len(options) >= 2 and answer:
            questions.append({"question": question, "options": options, "answer": answer})
    return questions


def section_label(sym: str, text: str):
    st.markdown(
        f'<div class="sect-label">'
        f'<span style="color:#86DDE4;font-size:0.9rem;line-height:1;">{sym}</span>'
        f'{text}</div>',
        unsafe_allow_html=True,
    )


def render_card(ph, text: str, streaming: bool = False,
                empty_sym: str = "✦", empty_msg: str = "Will appear here"):
    if text:
        if streaming:
            safe = _html.escape(text).replace("\n", "<br>")
            ph.markdown(
                f'<div class="output-card">{safe}<span class="cursor">|</span></div>',
                unsafe_allow_html=True,
            )
        else:
            rendered = _md.markdown(text, extensions=["nl2br", "tables"])
            ph.markdown(
                f'<div class="output-card md-output">{rendered}</div>',
                unsafe_allow_html=True,
            )
    else:
        ph.markdown(
            f'<div class="output-card empty">'
            f'<span style="font-size:1.6rem;color:rgba(196,53,86,0.4);">{empty_sym}</span>'
            f'<span>{empty_msg}</span></div>',
            unsafe_allow_html=True,
        )


def handle_error(e: Exception, model: str):
    err = str(e)
    if "rate limit" in err.lower() or "429" in err:
        st.error("Rate limit reached — wait a moment and try again.")
    elif "401" in err or "unauthorized" in err.lower():
        st.error("Invalid HF token — check your secrets.")
    elif "not a chat model" in err.lower() or "model_not_supported" in err:
        st.error(f"'{model}' does not support chat. Try a different model.")
    else:
        st.error(f"Something went wrong: {err}")


# ── Token check ───────────────────────────────────────────────────────────────
hf_token = get_hf_token()
if not hf_token:
    st.error("No HF token found. Add HF_TOKEN to .streamlit/secrets.toml")
    st.stop()

# ── Session state ─────────────────────────────────────────────────────────────
for k, v in {
    "fmt": "bullets", "last_summary": "", "last_terms": "",
    "history": [], "last_file_id": "", "last_image": None,
    "quiz_questions": [], "quiz_answers": {}, "quiz_submitted": False,
    "chat_history": [],
}.items():
    if k not in st.session_state:
        st.session_state[k] = v

# ── Sidebar ───────────────────────────────────────────────────────────────────
def _slabel(sym: str, text: str, mt: str = "0"):
    st.markdown(
        f'<div class="sidebar-label" style="margin-top:{mt};">'
        f'<span style="color:#86DDE4;font-size:0.9rem;line-height:1;">{sym}</span>'
        f'{text}</div>',
        unsafe_allow_html=True,
    )

with st.sidebar:
    # ── Header ──
    st.markdown("""
<div style="padding:1.5rem 0.5rem 0.4rem;text-align:center;position:relative;">
  <div style="position:absolute;top:0;left:50%;transform:translateX(-50%);
              width:52px;height:2px;border-radius:1px;
              background:linear-gradient(90deg,transparent,#C43556,transparent);"></div>
  <svg width="36" height="36" viewBox="0 0 24 24" fill="none"
       style="display:block;margin:0 auto 10px;">
    <path d="M12 3L2 9l10 6 10-6-10-6z"
          stroke="#86DDE4" stroke-width="1.6" stroke-linejoin="round" fill="rgba(134,221,228,0.08)"/>
    <path d="M6 12v5c0 2.5 2.7 4 6 4s6-1.5 6-4v-5"
          stroke="#86DDE4" stroke-width="1.6" stroke-linecap="round"/>
    <path d="M2 9v6" stroke="#C43556" stroke-width="1.6" stroke-linecap="round"/>
  </svg>
  <div style="font-family:'Playfair Display',serif;font-weight:700;font-size:1.2rem;
              color:#F9CCCB;letter-spacing:0.3px;">Caelune</div>
  <div style="font-size:0.67rem;letter-spacing:2px;text-transform:uppercase;
              color:rgba(134,221,228,0.5);margin-top:4px;">Study Assistant</div>
  <div style="margin:12px auto 0;width:40px;height:1px;
              background:linear-gradient(90deg,transparent,rgba(196,53,86,0.55),transparent);"></div>
</div>
""", unsafe_allow_html=True)

    # ── Output format ──
    _slabel("◈", "Output Format", mt="4px")
    fmt = st.radio(
        "fmt_r", label_visibility="collapsed", options=list(FORMAT_LABELS.keys()),
        format_func=lambda k: {"bullets": "Bullet Points", "flashcards": "Flashcards Q&A",
                               "structured": "Structured Guide"}[k],
        index=list(FORMAT_LABELS.keys()).index(st.session_state.fmt),
    )
    st.session_state.fmt = fmt
    st.markdown('<hr>', unsafe_allow_html=True)

    # ── Settings ──
    _slabel("○", "Subject")
    subject = st.selectbox("subj_s", label_visibility="collapsed", options=list(SUBJECTS.keys()))

    _slabel("△", "Difficulty", mt="10px")
    difficulty = st.radio("diff_r", label_visibility="collapsed",
                          options=list(DIFFICULTIES.keys()), horizontal=True)

    _slabel("≡", "Length", mt="10px")
    length = st.radio("len_r", label_visibility="collapsed",
                      options=list(LENGTHS.keys()), horizontal=True, index=1)

    _slabel("◐", "Language", mt="10px")
    language = st.selectbox("lang_s", label_visibility="collapsed", options=list(LANGUAGES.keys()))
    st.markdown('<hr>', unsafe_allow_html=True)

    _slabel("◇", "Model")
    model_label = st.selectbox("model_s", label_visibility="collapsed",
                               options=list(AVAILABLE_MODELS.keys()))
    selected_model = AVAILABLE_MODELS[model_label]
    st.markdown('<hr>', unsafe_allow_html=True)

    # ── History ──
    if st.session_state.history:
        _slabel("↺", "History")
        for i, entry in enumerate(reversed(st.session_state.history[-15:])):
            with st.expander(entry["title"], expanded=False):
                st.caption(f"{entry['fmt']} · {entry['subject']} · {entry['timestamp']}")
                if st.button("Restore", key=f"hist_{i}"):
                    st.session_state.last_summary = entry["summary"]
                    st.session_state.last_terms   = entry.get("terms", "")
                    st.rerun()
        st.markdown('<hr>', unsafe_allow_html=True)

    # ── Footer info ──
    st.markdown("""
<div class="info-row"><span style="color:#86DDE4;">◦</span>Hosted on Hugging Face</div>
<div class="info-row"><span style="color:#86DDE4;">◦</span>Notes are not stored</div>
<div class="info-row"><span style="color:#86DDE4;">◦</span>Free serverless inference</div>
""", unsafe_allow_html=True)

# ── Header ────────────────────────────────────────────────────────────────────
st.markdown("""
<div style="text-align:center;padding:0.5rem 0 1.6rem;">
  <h1>Caelune</h1>
  <p style="font-family:'DM Sans',sans-serif;font-size:0.93rem;color:rgba(134,221,228,0.55);font-weight:400;margin-top:5px;">
    Transform lecture notes into exam-ready summaries
  </p>
</div>
""", unsafe_allow_html=True)

# ── Layout ────────────────────────────────────────────────────────────────────
col_left, col_right = st.columns(2, gap="large")

with col_left:
    section_label("◉", "Your Notes")

    uploaded = st.file_uploader(
        "upload", label_visibility="collapsed",
        type=["pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg", "webp"],
        help="Upload a PDF, Word, PowerPoint, plain text, or an image (PNG/JPG/WEBP) of notes",
    )
    if uploaded is not None:
        file_id = f"{uploaded.name}_{uploaded.size}"
        if file_id != st.session_state.last_file_id:
            st.session_state.last_file_id = file_id
            ext = "." + uploaded.name.lower().rsplit(".", 1)[-1]
            if ext in _IMAGE_EXTS:
                raw = uploaded.read()
                st.session_state["last_image"] = raw
                with st.spinner("Extracting text from image…"):
                    extracted = extract_text_from_image(raw, ext, hf_token)
                if extracted:
                    st.session_state["notes_ta"] = extracted
            else:
                st.session_state["last_image"] = None
                st.session_state["notes_ta"] = extract_text_from_file(uploaded)
    else:
        st.session_state["last_image"] = None

    if st.session_state.get("last_image"):
        st.image(io.BytesIO(st.session_state["last_image"]),
                 caption="Uploaded image", use_container_width=True)

    notes = st.text_area(
        "notes_ta", key="notes_ta", label_visibility="collapsed",
        placeholder="Paste your lecture notes here, or upload a file above...",
        height=360,
    )
    char_count = len(notes)
    over = char_count > 50_000
    word_count = len(notes.split()) if notes.strip() else 0
    st.caption(f"{char_count:,} chars · {word_count:,} words{'  —  trim to continue' if over else ''}")
    st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)
    btn_col, clear_col = st.columns([4, 1])
    with btn_col:
        go = st.button("Generate Summary", type="primary", disabled=not notes.strip() or over)
    with clear_col:
        if st.button("Clear", type="secondary"):
            st.session_state["notes_ta"] = ""
            st.session_state["last_file_id"] = ""
            st.session_state["last_image"] = None
            st.rerun()

with col_right:
    section_label("✦", "Exam Summary")
    tab_summary, tab_terms, tab_quiz, tab_chat = st.tabs(["Summary", "Key Terms", "Quiz", "Chat"])

    with tab_summary:
        summary_ph = st.empty()
        render_card(summary_ph, st.session_state.last_summary,
                    empty_sym="✦", empty_msg="Your summary will appear here")
        if st.session_state.last_summary:
            st.download_button(
                "⬇ Download Summary",
                data=st.session_state.last_summary,
                file_name="caelune_summary.txt",
                mime="text/plain",
                key="dl_summary_persistent",
            )
    with tab_terms:
        terms_ph = st.empty()
        render_card(terms_ph, st.session_state.last_terms,
                    empty_sym="◈",
                    empty_msg="Key terms appear here after generating a summary")
        if st.session_state.last_terms:
            st.download_button(
                "⬇ Download Key Terms",
                data=st.session_state.last_terms,
                file_name="caelune_key_terms.txt",
                mime="text/plain",
                key="dl_terms_persistent",
            )

    # ── Quiz tab ──────────────────────────────────────────────────────────────
    with tab_quiz:
        if not notes.strip():
            st.markdown(
                '<div class="output-card empty">'
                '<span style="font-size:1.6rem;color:rgba(196,53,86,0.4);">△</span>'
                '<span>Enter or upload notes first, then generate a quiz</span></div>',
                unsafe_allow_html=True,
            )
        else:
            qn_col, qbtn_col = st.columns([3, 2])
            with qn_col:
                n_q = st.slider("Number of questions", 3, 10, 7, key="quiz_n_slider")
            with qbtn_col:
                st.markdown('<div style="height:27px"></div>', unsafe_allow_html=True)
                gen_quiz_btn = st.button("Generate Quiz", key="gen_quiz_btn", type="primary")

            if gen_quiz_btn:
                with st.spinner("Generating quiz questions…"):
                    raw_quiz = ""
                    try:
                        msgs = build_quiz_messages(
                            notes, n_q,
                            subject=subject, difficulty=difficulty, language=language,
                        )
                        for chunk in stream_call(msgs, hf_token, selected_model, 1400):
                            raw_quiz += chunk
                        parsed = parse_quiz(raw_quiz)
                        if parsed:
                            st.session_state.quiz_questions = parsed
                            st.session_state.quiz_submitted = False
                            st.session_state.quiz_answers = {}
                        else:
                            st.warning("Could not parse quiz output — try a different model or regenerate.")
                    except Exception as e:
                        handle_error(e, selected_model)

            questions = st.session_state.quiz_questions
            if questions:
                if not st.session_state.quiz_submitted:
                    st.markdown('<div style="height:6px"></div>', unsafe_allow_html=True)
                    for i, q in enumerate(questions):
                        st.markdown(
                            f'<p style="color:#F9CCCB;font-weight:600;font-size:0.9rem;'
                            f'margin:1rem 0 0.15rem;">'
                            f'Q{i + 1}. {_html.escape(q["question"])}</p>',
                            unsafe_allow_html=True,
                        )
                        st.radio(
                            f"q{i + 1}",
                            options=list(q["options"].keys()),
                            format_func=lambda k, opts=q["options"]: f"{k})  {opts[k]}",
                            key=f"quiz_q_{i}",
                            index=None,
                            label_visibility="collapsed",
                        )

                    answered = all(
                        st.session_state.get(f"quiz_q_{i}") is not None
                        for i in range(len(questions))
                    )
                    if not answered:
                        st.caption("Answer all questions to submit.")
                    if st.button("Check Answers", key="check_quiz_btn", type="primary",
                                 disabled=not answered):
                        st.session_state.quiz_answers = {
                            i: st.session_state.get(f"quiz_q_{i}")
                            for i in range(len(questions))
                        }
                        st.session_state.quiz_submitted = True
                        st.rerun()

                else:
                    answers = st.session_state.quiz_answers
                    score = sum(
                        1 for i, q in enumerate(questions)
                        if answers.get(i) == q["answer"]
                    )
                    pct = round(score / len(questions) * 100) if questions else 0
                    grade = "Excellent!" if pct >= 80 else ("Good effort!" if pct >= 60 else "Keep studying!")

                    st.markdown(f'''
<div class="quiz-score">
  <div class="quiz-score-num">{score} / {len(questions)}</div>
  <div class="quiz-score-label">{pct}% &nbsp;·&nbsp; {grade}</div>
</div>''', unsafe_allow_html=True)

                    results_html = ['<div style="max-height:340px;overflow-y:auto;padding-right:4px;">']
                    for i, q in enumerate(questions):
                        user_ans = answers.get(i)
                        correct = user_ans == q["answer"]
                        card_cls = "quiz-correct" if correct else "quiz-wrong"
                        badge_txt = "Correct ✓" if correct else "Wrong ✗"
                        badge_cls = "correct" if correct else "wrong"

                        opts_html = ""
                        for letter, text in q["options"].items():
                            is_right = letter == q["answer"]
                            is_user = letter == user_ans
                            if is_right:
                                style = "color:#86DDE4;font-weight:600;"
                                pfx = "✓ "
                            elif is_user:
                                style = "color:rgba(249,100,100,0.85);text-decoration:line-through;"
                                pfx = "✗ "
                            else:
                                style = "opacity:0.5;"
                                pfx = "&nbsp;&nbsp;&nbsp;"
                            opts_html += (
                                f'<div style="{style}padding:1px 0;font-size:0.86rem;">'
                                f'{pfx}<b>{letter})</b> {_html.escape(text)}</div>'
                            )

                        results_html.append(f'''
<div class="quiz-q {card_cls}" style="margin-bottom:0.7rem;">
  <span class="quiz-badge {badge_cls}">{badge_txt}</span>
  <div class="quiz-q-text" style="margin-bottom:0.45rem;">Q{i + 1}. {_html.escape(q["question"])}</div>
  {opts_html}
</div>''')
                    results_html.append('</div>')
                    st.markdown("".join(results_html), unsafe_allow_html=True)

                    rc1, rc2 = st.columns(2)
                    with rc1:
                        if st.button("Retake Quiz", key="retake_quiz_btn", type="secondary"):
                            st.session_state.quiz_submitted = False
                            st.session_state.quiz_answers = {}
                            st.rerun()
                    with rc2:
                        if st.button("New Quiz", key="new_quiz_btn", type="secondary"):
                            st.session_state.quiz_questions = []
                            st.session_state.quiz_submitted = False
                            st.session_state.quiz_answers = {}
                            st.rerun()

    # ── Chat tab ──────────────────────────────────────────────────────────────
    with tab_chat:
        if not notes.strip():
            st.markdown(
                '<div class="output-card empty">'
                '<span style="font-size:1.6rem;color:rgba(196,53,86,0.4);">◐</span>'
                '<span>Enter or upload notes first, then ask questions about them</span></div>',
                unsafe_allow_html=True,
            )
        else:
            if st.session_state.chat_history:
                bubbles = ['<div class="chat-container">']
                for msg in st.session_state.chat_history:
                    who = "You" if msg["role"] == "user" else "Caelune"
                    cls = "user" if msg["role"] == "user" else "ai"
                    if msg["role"] == "user":
                        content = _html.escape(msg["content"]).replace("\n", "<br>")
                    else:
                        content = _md.markdown(msg["content"], extensions=["nl2br", "tables"])
                    bubbles.append(
                        f'<div class="chat-bubble {cls}">'
                        f'<div class="chat-who">{who}</div>'
                        f'{content}</div>'
                    )
                bubbles.append('</div>')
                st.markdown("".join(bubbles), unsafe_allow_html=True)
            else:
                st.markdown(
                    '<p style="color:rgba(249,204,203,0.35);font-style:italic;'
                    'font-size:0.88rem;text-align:center;padding:1.2rem 0;">'
                    'Ask anything about your notes…</p>',
                    unsafe_allow_html=True,
                )

            st.markdown('<div style="height:6px"></div>', unsafe_allow_html=True)
            with st.form("chat_form", clear_on_submit=True):
                chat_q = st.text_input(
                    "chat_q", label_visibility="collapsed",
                    placeholder="Ask a question about your notes…",
                )
                send_btn = st.form_submit_button("Send →", type="primary")

            if send_btn and chat_q.strip():
                user_msg = chat_q.strip()
                st.session_state.chat_history.append({"role": "user", "content": user_msg})
                history_so_far = st.session_state.chat_history[:-1]
                msgs = build_chat_messages(
                    notes, history_so_far, user_msg,
                    subject=subject, language=language,
                )
                response = ""
                with st.spinner("Thinking…"):
                    try:
                        for chunk in stream_call(msgs, hf_token, selected_model, 700):
                            response += chunk
                    except Exception as e:
                        handle_error(e, selected_model)
                        st.session_state.chat_history.pop()
                if response:
                    st.session_state.chat_history.append(
                        {"role": "assistant", "content": response}
                    )
                st.rerun()

            if st.session_state.chat_history:
                if st.button("Clear Chat", key="clear_chat_btn", type="secondary"):
                    st.session_state.chat_history = []
                    st.rerun()

# ── Run ───────────────────────────────────────────────────────────────────────
if go:
    summary_collected = ""
    terms_collected   = ""
    max_tok = max_tokens_for(length)

    messages       = build_messages(
        notes, st.session_state.fmt,
        subject=subject, difficulty=difficulty, length=length, language=language,
    )
    terms_messages = build_key_terms_messages(notes, subject=subject, language=language)

    with col_right:
        with st.spinner("Generating summary..."):
            try:
                for chunk in stream_call(messages, hf_token, selected_model, max_tok):
                    summary_collected += chunk
                    with tab_summary:
                        render_card(summary_ph, summary_collected, streaming=True)
                with tab_summary:
                    render_card(summary_ph, summary_collected)
                st.session_state.last_summary = summary_collected
            except Exception as e:
                handle_error(e, selected_model)
                st.stop()

        if summary_collected:
            with st.spinner("Extracting key terms..."):
                try:
                    for chunk in stream_call(terms_messages, hf_token, selected_model, 650):
                        terms_collected += chunk
                        with tab_terms:
                            render_card(terms_ph, terms_collected, streaming=True,
                                        empty_sym="◈", empty_msg="")
                    with tab_terms:
                        render_card(terms_ph, terms_collected,
                                    empty_sym="◈", empty_msg="")
                    st.session_state.last_terms = terms_collected
                except Exception:
                    pass

        if summary_collected:
            title = (notes[:45].replace("\n", " ").strip() or "Untitled") + ("…" if len(notes) > 45 else "")
            st.session_state.history.append({
                "title":      title,
                "fmt":        st.session_state.fmt,
                "subject":    subject,
                "difficulty": difficulty,
                "length":     length,
                "language":   language,
                "summary":    summary_collected,
                "terms":      terms_collected,
                "timestamp":  datetime.now().strftime("%H:%M"),
            })
            st.rerun()
